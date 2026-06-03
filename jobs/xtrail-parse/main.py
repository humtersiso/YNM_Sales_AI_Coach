#!/usr/bin/env python3
"""
訓練素材：解析 pending 的 PDF/PPT → BigQuery knowledge_units（text_chunk）
環境變數：BIGQUERY_PROJECT_ID, BIGQUERY_DATASET, TRAINING_MATERIALS_ROOT（或 XTRAIL_ICE_SOURCE_ROOT）
用法：python main.py [--limit=50] [--asset-id=UUID]
"""
from __future__ import annotations

import argparse
import hashlib
import json
import os
import re
import uuid
from datetime import datetime, timezone
from pathlib import Path

from google.cloud import bigquery

try:
    from pypdf import PdfReader
except ImportError:
    PdfReader = None  # type: ignore

try:
    from pptx import Presentation
except ImportError:
    Presentation = None  # type: ignore


def env(name: str, default: str = "") -> str:
    return os.environ.get(name, default).strip()


def content_hash(parts: list[str]) -> str:
    return hashlib.sha256("\n".join(parts).encode("utf-8")).hexdigest()[:32]


def normalize_text(value: str) -> str:
    if not value:
        return ""
    text = (
        value.replace("\r\n", "\n")
        .replace("\r", "\n")
        .replace("\u000b", "\n")
    )
    cleaned = []
    for ch in text:
        o = ord(ch)
        if ch in ("\n", "\t") or o >= 32 or ch == "\u00a0":
            cleaned.append(ch)
        else:
            cleaned.append(" ")
    out = "".join(cleaned).replace("\ufffd", "")
    while "\n\n\n" in out:
        out = out.replace("\n\n\n", "\n\n")
    return out.strip()


NOISE_LINE = re.compile(
    r"all rights reserved|confidential|yulon\s*nissan|do not use|permission|copyright|^\d+\s*$",
    re.I,
)
SPEC_TOKEN = re.compile(
    r"\d+(?:\.\d+)?\s*(?:ps|PS|kgm|kg·m|km/L|km/l|匹|公里/公升|公斤米|牛頓米|Nm)",
    re.I,
)


def extract_script_excerpt(script: str, max_chars: int = 380) -> str:
    lines = [ln.strip() for ln in script.replace("\r\n", "\n").split("\n") if len(ln.strip()) >= 2]
    meaningful: list[str] = []
    for line in lines:
        if NOISE_LINE.search(line):
            continue
        cjk = sum(1 for ch in line if "\u4e00" <= ch <= "\u9fff")
        if cjk < 2 and len(line) < 24:
            continue
        meaningful.append(line)
        if len(" ".join(meaningful)) >= max_chars or len(meaningful) >= 8:
            break
    out = " ".join(meaningful).strip()
    specs = [m.group(0).strip() for m in SPEC_TOKEN.finditer(script)][:6]
    for token in specs:
        if token.lower().replace(" ", "") not in out.lower().replace(" ", ""):
            out = f"{out} {token}".strip() if out else token
    return out[:max_chars]


def build_chunk_fields(file_name: str, locator_key: str, loc: int, script: str) -> tuple[str, str]:
    title = f"{file_name} ({locator_key} {loc})"
    excerpt = extract_script_excerpt(script)
    stem = re.sub(r"\.(pdf|pptx|ppt)$", "", file_name, flags=re.I).strip()
    if excerpt:
        cq = f"{stem} ({locator_key} {loc}) · {excerpt}".replace("  ", " ")[:500]
    else:
        cq = title[:500]
    return title[:300], cq


def is_garbled(text: str) -> bool:
    t = normalize_text(text)
    if len(t) < 4:
        return False
    if t.startswith("PK") and ("[Content_Types]" in t or "xmlschemas" in t):
        return True
    if "[Content_Types].xml" in t or "_rels/.rels" in t:
        return True
    if t.count("\ufffd") > 2:
        return True
    cjk = sum(1 for ch in t if "\u4e00" <= ch <= "\u9fff")
    printable = sum(
        1
        for ch in t
        if "\u4e00" <= ch <= "\u9fff"
        or ch.isalnum()
        or ch in " ，。、？！：；「」（）-_%"
        or ch.isspace()
    )
    if len(t) > 40 and cjk == 0 and printable / len(t) < 0.35:
        return True
    return False


def resolve_local_path(relative_path: str, root: Path, product_line: str) -> Path:
    rel = relative_path.replace("/", os.sep)
    direct = (root / rel).resolve()
    if direct.is_file():
        return direct
    nested = (root / product_line / rel).resolve()
    return nested


def extract_pdf_text(path: Path) -> list[tuple[int, str]]:
    if PdfReader is None:
        raise RuntimeError("pypdf 未安裝")
    reader = PdfReader(str(path))
    chunks: list[tuple[int, str]] = []
    for i, page in enumerate(reader.pages, start=1):
        text = normalize_text(page.extract_text() or "")
        if text and not is_garbled(text):
            chunks.append((i, text))
    return chunks


def extract_pptx_text(path: Path) -> list[tuple[int, str]]:
    if Presentation is None:
        raise RuntimeError("python-pptx 未安裝")
    prs = Presentation(str(path))
    chunks: list[tuple[int, str]] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text.strip())
        text = normalize_text("\n".join(p for p in parts if p))
        if text and not is_garbled(text):
            chunks.append((i, text))
    return chunks


def list_pending(client: bigquery.Client, project: str, dataset: str, limit: int, asset_id: str | None):
    table = f"`{project}.{dataset}.source_assets`"
    where = "parse_status = 'pending'"
    params: list[bigquery.ScalarQueryParameter] = []
    if asset_id:
        where += " AND asset_id = @asset_id"
        params.append(bigquery.ScalarQueryParameter("asset_id", "STRING", asset_id))
    sql = f"""
        SELECT *
        FROM {table}
        WHERE {where}
          AND (
            LOWER(file_name) LIKE '%.pdf'
            OR LOWER(file_name) LIKE '%.pptx'
            OR LOWER(file_name) LIKE '%.ppt'
          )
        ORDER BY ingested_at
        LIMIT @lim
    """
    params.append(bigquery.ScalarQueryParameter("lim", "INT64", limit))
    job_config = bigquery.QueryJobConfig(query_parameters=params)
    return list(client.query(sql, job_config=job_config).result())


def update_asset_status(
    client: bigquery.Client,
    project: str,
    dataset: str,
    asset_id: str,
    status: str,
    error: str | None,
):
    """streaming insert 後短時間內無法 UPDATE，失敗時略過（knowledge_units 仍有效）。"""
    sql = f"""
        UPDATE `{project}.{dataset}.source_assets`
        SET parse_status = @status, parse_error = @error
        WHERE asset_id = @asset_id
    """
    job_config = bigquery.QueryJobConfig(
        query_parameters=[
            bigquery.ScalarQueryParameter("status", "STRING", status),
            bigquery.ScalarQueryParameter("error", "STRING", error),
            bigquery.ScalarQueryParameter("asset_id", "STRING", asset_id),
        ]
    )
    try:
        client.query(sql, job_config=job_config).result()
    except Exception as e:
        if "streaming buffer" in str(e).lower():
            print(
                json.dumps(
                    {
                        "warn": "skip_asset_status_update",
                        "asset_id": asset_id,
                        "status": status,
                    },
                    ensure_ascii=False,
                )
            )
            return
        raise


def insert_units(client: bigquery.Client, project: str, dataset: str, rows: list[dict]):
    if not rows:
        return
    table_id = f"{project}.{dataset}.knowledge_units"
    errors = client.insert_rows_json(client.get_table(table_id), rows)
    if errors:
        raise RuntimeError(f"insert_rows_json failed: {errors[:3]}")


def tags_from_path(relative_path: str, product_line: str) -> list[str]:
    parts = relative_path.replace("\\", "/").split("/")
    if len(parts) <= 1:
        return [product_line] if product_line else []
    tags = [product_line] + [p.strip() for p in parts[:-1] if p.strip() and len(p) < 80]
    return list(dict.fromkeys(tags))[:10]


def process_asset(client, project, dataset, asset, root: Path, dry_run: bool) -> int:
    rel = asset["relative_path"]
    product_line = asset.get("product_line") or "_common"
    material_category = asset.get("material_category") or "general"
    local = resolve_local_path(rel, root, product_line)
    if not local.is_file():
        update_asset_status(client, project, dataset, asset["asset_id"], "failed", f"找不到檔案: {local}")
        return 0

    name_lower = asset["file_name"].lower()
    try:
        if name_lower.endswith(".pdf"):
            page_chunks = extract_pdf_text(local)
            locator_key = "page"
        elif name_lower.endswith(".pptx") or name_lower.endswith(".ppt"):
            page_chunks = extract_pptx_text(local)
            locator_key = "slide"
        else:
            update_asset_status(client, project, dataset, asset["asset_id"], "unsupported", "非 PDF/PPT")
            return 0
    except Exception as e:
        update_asset_status(client, project, dataset, asset["asset_id"], "failed", str(e))
        return 0

    if not page_chunks:
        update_asset_status(client, project, dataset, asset["asset_id"], "failed", "無可抽取文字")
        return 0

    now = datetime.now(timezone.utc).isoformat()
    tags = tags_from_path(rel, product_line)
    units: list[dict] = []
    for loc, text in page_chunks:
        locator = json.dumps({locator_key: loc}, ensure_ascii=False)
        script = text[:15000]
        title, question = build_chunk_fields(asset["file_name"], locator_key, loc, script)
        h = content_hash([product_line, material_category, asset["asset_id"], str(loc), script[:500]])
        units.append(
            {
                "unit_id": str(uuid.uuid4()),
                "ingest_batch_id": asset["ingest_batch_id"],
                "asset_id": asset["asset_id"],
                "product_line": product_line,
                "material_category": material_category,
                "unit_type": "text_chunk",
                "title": title,
                "customer_question": question,
                "standard_script": script,
                "source_locator": locator,
                "tags": tags,
                "language": "zh-TW",
                "content_hash": h,
                "ingested_at": now,
            }
        )

    if not dry_run:
        insert_units(client, project, dataset, units)
        update_asset_status(client, project, dataset, asset["asset_id"], "ok", None)
    return len(units)


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--limit", type=int, default=50)
    parser.add_argument("--asset-id", default="")
    parser.add_argument("--dry-run", action="store_true")
    args = parser.parse_args()

    project = env("BIGQUERY_PROJECT_ID", env("GOOGLE_CLOUD_PROJECT"))
    dataset = env("BIGQUERY_DATASET", "YNM_Sales_AI_Coach_test")
    default_root = Path(__file__).resolve().parents[2] / "data" / "training-materials"
    root = Path(
        env("TRAINING_MATERIALS_ROOT")
        or env("XTRAIL_ICE_SOURCE_ROOT")
        or str(default_root)
    )

    if not project:
        raise SystemExit("請設定 BIGQUERY_PROJECT_ID")

    client = bigquery.Client(project=project)
    assets = list_pending(
        client,
        project,
        dataset,
        args.limit,
        args.asset_id or None,
    )

    total_units = 0
    for asset in assets:
        total_units += process_asset(client, project, dataset, dict(asset.items()), root, args.dry_run)

    print(json.dumps({"processed": len(assets), "units_written": total_units, "root": str(root)}, ensure_ascii=False))


if __name__ == "__main__":
    main()
